import openpyxl as a
from openpyxl.utils import get_column_letter
import pptx as b
from pptx.chart.data import CategoryChartData

answers = []
theMostAns = []
contQuests = []
podContQuests = {}
tables = {}
row = 2
col = 1

def SetTables():
    for ans in answers:
        for quest in ans:
            if tables.get(quest) is None:
                tables[quest] = {}
            try:
                if tables[quest].get(ans[quest]['Ответ']) is None:
                    tables[quest][ans[quest].get('Ответ')] = 1
                else:
                    tables[quest][ans[quest].get('Ответ')] = tables[quest][ans[quest].get('Ответ')] + 1
            except:
                for variant in ans[quest]:
                    if ans[quest][variant] is not None:
                        if tables[quest].get(variant) is None:
                            tables[quest][variant] = 1
                        else:
                            tables[quest][variant] += 1

    for quest in tables:
        try:
            tables[quest] = dict(sorted(tables[quest].items()))
        except:
            pass

    # for i in tables:
    #     print(f'{i} {tables[i]}')


def GlobStat(pod, qVariants):
    global podContQuests
    podContQuests=pod
    wb = a.load_workbook("./open.xlsx")
    try:
        wb.remove(wb['частота'])
    except:
        pass
    sheet = wb.create_sheet('частота')
    sheet.column_dimensions['A'].width = 50
    sheet.column_dimensions['B'].width = 10
    sheet.column_dimensions['c'].width = 10

    stats = {}
    global col, row
    row = 2
    for quest in tables:
        col = 1
        sheet[f'{get_column_letter(col)}{row}'] = f'{quest}~Ответ'
        col += 1
        sheet[f'{get_column_letter(col)}{row}'] = 'кол-во'
        col += 1
        sheet[f'{get_column_letter(col)}{row}'] = '%'
        contRow=row

        isContQuest=False
        catNames=[]
        for i in podContQuests:
            if quest in podContQuests[i]:
                isContQuest = True
                catNames.append(i)
        summ=0
        specSumm=0
        for variant in tables[quest]:
            if variant is None:
                continue
            col = 1
            row += 1
            sheet[f'{get_column_letter(col)}{row}'] = variant
            col += 1
            sheet[f'{get_column_letter(col)}{row}'] = tables[quest][variant]
            summ += tables[quest][variant]
        row=contRow
        for variant in tables[quest]:
            if variant is None:
                continue
            col = 1
            row += 1
            col += 2
            per = (100 / summ * tables[quest][variant])
            sheet[f'{get_column_letter(col)}{row}'] = per
            if isContQuest:
                for catName in catNames:
                    if quest in podContQuests[catName] and variant.lower() in qVariants[catName]['values']:
                        if qVariants[catName]['values'][variant.lower()] == 'green':
                            specSumm+=per
                        elif not qVariants[catName]['type'] and qVariants[catName]['values'][variant.lower()] == 'red':
                            specSumm-=per
                        stats[catName] = (stats[catName] + per) if stats.get(catName) is not None else per
        if isContQuest:
            row+=1
            sheet[f'{get_column_letter(col)}{row}'] = specSumm
        row += 2
    # stats['Самые вовлеченные']=TheMost()
    extCol = sheet.max_column + 5
    row = 2
    for stat in podContQuests:
        col = extCol
        sheet[f'{get_column_letter(col)}{row}'] = stat
        col += 1
        sheet[f'{get_column_letter(col)}{row}'] = (stats[stat] / len(podContQuests[stat])) if stats.get(stat) is not None else 0
        row += 1

    wb.save('./open.xlsx')

def TheMost():
    theMost = 0
    i = 0
    for ans in answers:
        combo = 0
        for cat in podContQuests:
            for quest in cat:
                if ans[quest].get('Ответ') is not None and ans[quest]['Ответ'].lower() in ['да', 'скорее да']:
                    combo += 1
            if 12 == combo:
                theMost += 1
                global theMostAns
                theMostAns.append(ans)
            i += 1
    return theMost

def ByAllStats(statsQuests, qVariants):
    wb = a.load_workbook("./open.xlsx")
    for quest in statsQuests:
        try:
            wb.remove(wb[f'by {quest[0:10]}'])
        except:
            pass
        sheet = wb.create_sheet(f'by {quest[0:10]}')
        sheet.column_dimensions['A'].width = 50
        stats = {}
        global col, row
        row = 2
        for squest in tables:
            if quest == squest:
                continue
            col = 1
            sheet[f'{get_column_letter(col)}{row}'] = f'{quest}~Ответ'
            for i in tables[quest]:
                if i is None:
                    continue
                col += 1
                sheet[f'{get_column_letter(col)}{row}'] = i
                sheet.column_dimensions[get_column_letter(col)].width = 10
                col += 1
                sheet[f'{get_column_letter(col)}{row}'] = f'% {i}'
                sheet.column_dimensions[get_column_letter(col)].width = 10
            row += 1
            col = 1
            sheet[f'{get_column_letter(col)}{row}'] = f'{squest}~Ответ'
            row += 1
            nums = GetNum(quest, squest)
            isContQuest=False
            catNames=[]
            for i in podContQuests:
                if squest in podContQuests[i]:
                    isContQuest=True
                    catNames.append(i)
            specSumm={}
            for i in tables[squest]:
                if i is None:
                    continue
                sheet[f'{get_column_letter(col)}{row}'] = i
                for j in tables[quest]:
                    if j is None:
                        continue
                    col += 1
                    value=nums[j].get(i) if (nums.get(j) is not None and nums[j].get(i) is not None) else 0
                    sheet[f'{get_column_letter(col)}{row}'] = value
                    col += 1
                    per = (100 / nums[j]['sum'] * value) if nums[j]['sum'] != 0 else 0
                    sheet[f'{get_column_letter(col)}{row}'] = per
                    if isContQuest:
                        for catName in catNames:
                            if squest in podContQuests[catName] and i.lower() in qVariants[catName]['values']:
                                if qVariants[catName]['values'][i.lower()] == 'green':
                                    specSumm[j] = specSumm[j] + per if specSumm.get(j) is not None else per
                                elif not qVariants[catName]['type'] and qVariants[catName]['values'][i.lower()] == 'red':
                                    specSumm[j] = specSumm[j] + per if specSumm.get(j) is not None else per
                                if stats.get(j) is None:
                                    stats[j] = {}
                                stats[j][catName] = (stats[j][catName] + per if stats[j].get(catName) is not None else per)
                row += 1
                col = 1
            col=1
            if isContQuest:
                for i in tables[quest]:
                    col+=2
                    sheet[f'{get_column_letter(col)}{row}'] = specSumm[i] if specSumm.get(i) is not None else 0
                row+=1
            row += 1
        extCol = sheet.max_column + 5
        row = 2
        col = extCol
        for variant in tables[quest]:
            col += 1
            sheet[f'{get_column_letter(col)}{row}'] = variant
        row += 1
        for stat in podContQuests:
            col = extCol
            sheet[f'{get_column_letter(col)}{row}'] = stat
            for variant in tables[quest]:
                col += 1
                sheet[f'{get_column_letter(col)}{row}'] = (stats[variant][stat] / len(podContQuests[stat])) if stats.get(variant) is not None and stats[variant].get(
                    stat) is not None else 0
            row += 1
    wb.save('./open.xlsx')


def GetNum(quest1, quest2):
    gg = {}
    for ans in answers:
        for i in ans[quest1]:
            if ans[quest1][i] is not None:
                if gg.get(ans[quest1][i]) is None:
                    gg[ans[quest1][i]] = {}
                    gg[ans[quest1][i]]['sum'] = 0
                for j in ans[quest2]:
                    if ans[quest2][j] is not None:
                        if gg[ans[quest1][i]].get(ans[quest2][j]) is not None:
                            gg[ans[quest1][i]]['sum'] += 1
                            gg[ans[quest1][i]][ans[quest2][j]] += 1
                        else:
                            gg[ans[quest1][i]]['sum'] += 1
                            gg[ans[quest1][i]][ans[quest2][j]] = 1

    return gg


def FPages():
    prs = b.Presentation("shablon.pptx")

    # slide 1
    chart = prs.slides[0].shapes[6].chart
    SetChart(chart, 'возраст')

    chart = prs.slides[0].shapes[16].chart
    SetChart(chart, 'работ')

    GetPercent(prs.slides[0].shapes[18], prs.slides[0].shapes[17], 'пол')
    GetPercent(prs.slides[0].shapes[23], prs.slides[0].shapes[24], 'подчин')

    # slide 3
    # chart = prs.slides[2].shapes[24].chart
    # SetChart(chart, 'возраст', isTheMost=True)
    #
    # chart = prs.slides[2].shapes[25].chart
    # SetChart(chart, 'работ', isTheMost=True)
    #
    # GetPercent(prs.slides[2].shapes[15], prs.slides[2].shapes[14], 'пол', isTheMost=True)
    # GetPercent(prs.slides[2].shapes[20], prs.slides[2].shapes[21], 'подчин', isTheMost=True)

    prs.save("Statistics.pptx")


def GetPercent(value1, value2, key, isTheMost=False):
    cats = {}
    summ = 0
    if not isTheMost:
        for quest in tables:
            if quest.lower().find(key) != -1:
                for variant in tables[quest]:
                    cats[variant] = tables[quest][variant]
                    summ += tables[quest][variant]
                break
        if key == 'пол':
            for variant in cats:
                if variant.lower().find('муж'):
                    value1.text_frame.paragraphs[0].runs[0].text = f'{round(100 / summ * cats[variant])}%'
                if variant.lower().find('жен'):
                    value2.text_frame.paragraphs[0].runs[0].text = f'{round(100 / summ * cats[variant])}%'
        elif key == 'подчин':
            for variant in cats:
                if variant.lower().find('да'):
                    value1.text_frame.paragraphs[0].runs[0].text = f'{round(100 / summ * cats[variant])}%'
                if variant.lower().find('нет'):
                    value2.text_frame.paragraphs[0].runs[0].text = f'{round(100 / summ * cats[variant])}%'
    else:
        for ans in theMostAns:
            for quest in ans:
                if quest.lower().find(key) != -1:
                    for variant in ans[quest]:
                        cats[ans[quest][variant]] = cats[ans[quest][variant]] + 1 if cats.get(
                            ans[quest][variant]) is not None else 1
                    break
        summ = {}
        for ans in answers:
            for quest in ans:
                if quest.lower().find(key) != -1:
                    for variant in ans[quest]:
                        summ[ans[quest][variant]] = summ[ans[quest][variant]] + 1 if summ.get(
                            ans[quest][variant]) is not None else 1
                    break
        if key == 'пол':
            for variant in cats:
                if variant.lower().find('муж'):
                    value1.text_frame.paragraphs[0].runs[0].text = f'{round(100 / summ[variant] * cats[variant])}%'
                if variant.lower().find('жен'):
                    value2.text_frame.paragraphs[0].runs[0].text = f'{round(100 / summ[variant] * cats[variant])}%'
        elif key == 'подчин':
            for variant in cats:
                if variant.lower().find('да'):
                    value1.text_frame.paragraphs[0].runs[0].text = f'{round(100 / summ[variant] * cats[variant])}%'
                if variant.lower().find('нет'):
                    value2.text_frame.paragraphs[0].runs[0].text = f'{round(100 / summ[variant] * cats[variant])}%'


def SetChart(chart, key, isTheMost=False):
    chart_data = CategoryChartData()
    cats = {}
    summ = 0
    if not isTheMost:
        for quest in tables:
            if quest.lower().find(key) != -1:
                for variant in tables[quest]:
                    cats[variant] = tables[quest][variant]
                    summ += tables[quest][variant]
                break
        values = []
        for variant in cats:
            values.append(round(100 / summ * cats[variant]))
    else:
        for ans in theMostAns:
            for quest in ans:
                if quest.lower().find(key) != -1:
                    for variant in ans[quest]:
                        cats[ans[quest][variant]] = cats[ans[quest][variant]] + 1 if cats.get(
                            ans[quest][variant]) is not None else 1
                    break
        summ = {}
        for ans in answers:
            for quest in ans:
                if quest.lower().find(key) != -1:
                    for variant in ans[quest]:
                        summ[ans[quest][variant]] = summ[ans[quest][variant]] + 1 if summ.get(
                            ans[quest][variant]) is not None else 1
                    break
        values = []
        for variant in cats:
            values.append(round(100 / summ[variant] * cats[variant]))
    chart_data.categories = cats.keys()
    chart_data.add_series('Series 1', values)
    chart.replace_data(chart_data)
    pass
